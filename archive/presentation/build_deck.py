"""Build the STAT 4830 presentation deck (python-pptx).

Editorial theme: paper background, serif display type, hairline rules,
per-track color coding (Track 1 navy, Track 2 forest, Track 3 brick).

13 slides, 10-15 minutes, smart grad audience.
Three tracks: cascade GPU (1.4), Lasserre SDP (1.3 certified), MATLAB audit.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "figures"
OUT = HERE.parent / "STAT 4830 Presentation.pptx"

# --- Editorial palette (matches generate_figures.py) ---
PAPER = RGBColor(0xF7, 0xF2, 0xE8)
INK = RGBColor(0x1C, 0x19, 0x17)
STONE = RGBColor(0x7C, 0x75, 0x6B)
HAIRLINE = RGBColor(0xD6, 0xCD, 0xB8)
HIGHLIGHT = RGBColor(0xEE, 0xE6, 0xD2)

NAVY = RGBColor(0x25, 0x43, 0x5F)     # Track 1 — cascade
FOREST = RGBColor(0x3A, 0x54, 0x45)   # Track 2 — SDP
BRICK = RGBColor(0x8E, 0x3A, 0x38)    # Track 3 — audit
TERRA = RGBColor(0xB5, 0x4B, 0x2F)    # cross-section accent
GOLD = RGBColor(0xA8, 0x84, 0x31)     # secondary accent

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# HELPERS
# ============================================================

def new_prs() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank_slide(prs, *, rail_color=None):
    """Create a blank slide with the editorial paper background.

    If rail_color is given, draws a thin vertical color rule on the left edge.
    """
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    # paper background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    # optional left rail
    if rail_color is not None:
        rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  Inches(0.18), SLIDE_H)
        rail.fill.solid()
        rail.fill.fore_color.rgb = rail_color
        rail.line.fill.background()
        rail.shadow.inherit = False
    return s


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=SANS, spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_rich(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.35):
    """runs = list of paragraphs; each paragraph = list of (text, dict)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for text, style in para:
            r = p.add_run()
            r.text = text
            r.font.name = style.get("font", SANS)
            r.font.size = Pt(style.get("size", 18))
            r.font.bold = style.get("bold", False)
            r.font.italic = style.get("italic", False)
            r.font.color.rgb = style.get("color", INK)
    return tb


def add_rule(slide, left, top, width, color=HAIRLINE, thickness=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_box(slide, left, top, width, height, fill=HIGHLIGHT, line=None,
            line_thickness=0.6):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_thickness)
    sh.shadow.inherit = False
    return sh


def add_dot(slide, cx, cy, radius, color):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - radius, cy - radius, radius * 2, radius * 2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def slide_header(slide, eyebrow, title, *, color=INK, eyebrow_color=None):
    """Editorial header: small eyebrow label, big serif title, hairline rule."""
    if eyebrow_color is None:
        eyebrow_color = color if color is not INK else STONE
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.35),
                 eyebrow.upper(),
                 size=11, bold=True, color=eyebrow_color,
                 font=SANS)
    add_text(slide, Inches(0.6), Inches(0.72), Inches(12.1), Inches(0.8),
             title, size=30, bold=False, color=INK, font=SERIF)
    add_rule(slide, Inches(0.6), Inches(1.46), Inches(12.1),
             color=HAIRLINE, thickness=0.75)


def add_footer(slide, idx, total, *, active_color=INK):
    """Footer: authors / course on the left; progress dots on the right."""
    add_text(slide, Inches(0.6), Inches(7.12), Inches(8),
             Inches(0.3),
             "Piterbarg  ·  Bajaj  ·  Vincent        STAT 4830  —  Spring 2026",
             size=9, color=STONE, font=SANS)
    # progress dots — active dot in slide's accent color
    right = Inches(12.73)
    gap = Inches(0.18)
    radius = Inches(0.045)
    # draw from right to left so right edge anchors to `right`
    for i in range(total, 0, -1):
        cx = right - (total - i) * gap
        cy = Inches(7.22)
        color = active_color if i == idx else HAIRLINE
        add_dot(slide, cx, cy, radius, color)


# ============================================================
# SLIDES
# ============================================================

def slide_01_title(prs):
    s = blank_slide(prs)

    # Top eyebrow
    add_text(s, Inches(0.8), Inches(0.8), Inches(11.8), Inches(0.4),
             "STAT 4830    ·    SPRING 2026",
             size=11, bold=True, color=STONE, font=SANS)
    add_rule(s, Inches(0.8), Inches(1.25), Inches(2.6),
             color=INK, thickness=1.1)

    # Big serif title
    add_text(s, Inches(0.8), Inches(2.05), Inches(11.8), Inches(1.0),
             "Lower Bounds for the",
             size=52, bold=False, color=INK, font=SERIF, spacing=1.02)
    add_text(s, Inches(0.8), Inches(2.9), Inches(11.8), Inches(1.2),
             "Autoconvolution Constant",
             size=52, bold=False, italic=True, color=INK, font=SERIF,
             spacing=1.02)

    # Subtitle
    add_rich(s, Inches(0.8), Inches(4.15), Inches(11.8), Inches(0.6), [
        [("Three parallel attacks on   ", {"size": 22, "color": STONE}),
         ("C₁ₐ", {"size": 26, "color": INK, "italic": True, "font": SERIF})],
    ])

    # hairline divider
    add_rule(s, Inches(0.8), Inches(5.25), Inches(11.8),
             color=HAIRLINE, thickness=0.75)

    # byline
    add_text(s, Inches(0.8), Inches(5.45), Inches(11.8), Inches(0.5),
             "Andrei Piterbarg     Jai Bajaj     Derrick Vincent",
             size=18, bold=False, color=INK, font=SANS)

    # track color swatches as a signature
    swatch_y = Inches(6.35)
    swatch_w = Inches(0.5)
    swatch_h = Inches(0.12)
    swatch_gap = Inches(0.18)
    start = Inches(0.8)
    for i, (label, col) in enumerate([
        ("I.  Cascade", NAVY),
        ("II.  SDP certificate", FOREST),
        ("III.  MATLAB audit", BRICK),
    ]):
        x = start + i * (swatch_w + swatch_gap + Inches(2.2))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, swatch_y,
                                 swatch_w, swatch_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        bar.shadow.inherit = False
        add_text(s, x + swatch_w + Inches(0.12), swatch_y - Inches(0.04),
                 Inches(2.5), Inches(0.3),
                 label, size=12, bold=True, color=INK, font=SANS)


def slide_02_problem(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "§1 · The problem",
                 "How small can the peak of a normalized self-convolution be?")

    add_rich(s, Inches(0.6), Inches(1.85), Inches(6.3), Inches(3.5), [
        [("Given  ", {"size": 18}),
         ("f : ℝ → ℝ₊", {"size": 19, "italic": True, "color": NAVY,
                        "font": SERIF}),
         ("  with  supp ", {"size": 18}),
         ("f", {"size": 18, "italic": True, "font": SERIF}),
         ("  ⊆  [−¼, ¼]", {"size": 18})],
        [("and  ", {"size": 18}),
         ("∫ f = 1", {"size": 18, "font": SERIF, "italic": True}),
         (",   define", {"size": 18})],
        [("", {"size": 16})],
        [("C₁ₐ  =  inf  ‖ f ∗ f ‖∞",
          {"size": 34, "bold": False, "color": INK, "font": SERIF,
           "italic": True})],
        [("", {"size": 18})],
        [("·  connected to Sidon  /  generalized Sidon set densities",
          {"size": 14, "color": STONE})],
        [("·  best published bounds  have not moved  since 2017",
          {"size": 14, "color": STONE})],
        [("·  every improvement is a real contribution",
          {"size": 14, "color": STONE})],
    ])

    s.shapes.add_picture(str(FIG / "fig_autoconv_intuition.png"),
                         Inches(7.05), Inches(2.0), width=Inches(5.9))

    # Bottom bounds line — minimal, typographic
    add_rule(s, Inches(0.6), Inches(5.7), Inches(12.1),
             color=HAIRLINE, thickness=0.75)
    add_text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.3),
             "BEST PUBLISHED BOUNDS",
             size=10, bold=True, color=STONE, font=SANS)
    add_rich(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.6), [
        [("1.2802", {"size": 28, "bold": True, "color": NAVY,
                     "font": SERIF}),
         ("   ≤   ", {"size": 24, "color": STONE}),
         ("C₁ₐ", {"size": 28, "italic": True, "font": SERIF}),
         ("   ≤   ", {"size": 24, "color": STONE}),
         ("1.5029", {"size": 28, "bold": True, "color": BRICK,
                     "font": SERIF}),
         ("        unchanged since 2017",
          {"size": 14, "color": STONE, "italic": True})],
    ])

    add_footer(s, idx, total, active_color=INK)


def slide_03_baseline_and_tracks(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "§2 · Baseline and our three tracks",
                 "Three independent attacks above the 2017 floor")

    # PRIOR WORK row (left summary + right result card) ------------
    add_text(s, Inches(0.6), Inches(1.75), Inches(8.0), Inches(0.35),
             "PRIOR WORK  —  CLONINGER & STEINERBERGER 2017",
             size=11, bold=True, color=STONE, font=SANS)
    add_rich(s, Inches(0.6), Inches(2.1), Inches(8.0), Inches(1.0), [
        [("Branch-and-prune over the fine height grid  ",
          {"size": 14, "color": INK}),
         ("B", {"size": 14, "italic": True, "font": SERIF}),
         ("n,m", {"size": 10, "italic": True, "font": SERIF}),
         (":   discretize  [−¼, ¼]  into  ", {"size": 14, "color": INK}),
         ("d = 2n", {"size": 14, "italic": True, "font": SERIF}),
         ("  bins,", {"size": 14, "color": INK})],
        [("enumerate height vectors and prune any whose windowed test value exceeds the target;  correction term  ",
          {"size": 14, "color": INK}),
         ("2/m + 1/m²", {"size": 14, "italic": True, "font": SERIF}),
         (".", {"size": 14, "color": INK})],
    ], spacing=1.3)

    # Right: compact result card
    add_box(s, Inches(8.9), Inches(1.75), Inches(3.8), Inches(1.45),
            fill=HIGHLIGHT)
    add_rule(s, Inches(8.9), Inches(1.75), Inches(3.8),
             color=INK, thickness=1.3)
    add_text(s, Inches(9.05), Inches(1.87), Inches(3.5), Inches(0.3),
             "PUBLISHED RESULT",
             size=10, bold=True, color=STONE, font=SANS)
    add_text(s, Inches(9.05), Inches(2.2), Inches(3.5), Inches(0.6),
             "C₁ₐ  ≥  1.2802",
             size=26, bold=True, color=INK, font=SERIF)
    add_text(s, Inches(9.05), Inches(2.8), Inches(3.5), Inches(0.3),
             "(n, m) = (24, 50)   ·   ≈ 20,000 CPU-hours",
             size=11, color=STONE, italic=True, font=SANS)

    # Divider between prior work and our tracks
    add_rule(s, Inches(0.6), Inches(3.4), Inches(12.1),
             color=HAIRLINE, thickness=0.75)

    # OUR CONTRIBUTION row --------------------------------------------
    add_text(s, Inches(0.6), Inches(3.55), Inches(12.1), Inches(0.35),
             "OUR CONTRIBUTION  —  THREE INDEPENDENT TRACKS",
             size=11, bold=True, color=STONE, font=SANS)

    s.shapes.add_picture(str(FIG / "fig_three_tracks.png"),
                         Inches(1.67), Inches(3.95), width=Inches(10.0))

    add_footer(s, idx, total, active_color=INK)


def slide_05_cascade_method(prs, idx, total):
    s = blank_slide(prs, rail_color=NAVY)
    slide_header(s, "TRACK I · CASCADE",
                 "Recursive refinement with novel pruning, fused GPU kernels",
                 color=NAVY, eyebrow_color=NAVY)

    add_rich(s, Inches(0.6), Inches(1.8), Inches(6.7), Inches(5.0), [
        [("Novel pruning bounds",
          {"size": 17, "bold": True, "color": NAVY, "font": SERIF})],
        [("·  contiguous-block sums catch interior concentrations",
          {"size": 14})],
        [("·  two-max enhanced bound — cross-term dominance", {"size": 14})],
        [("·  edge-pair cross-terms  ", {"size": 14}),
         ("c₀ · c_{d−1}", {"size": 14, "font": SERIF, "italic": True})],
        [("each bound proved, then compiled into the pruner",
          {"size": 12, "italic": True, "color": STONE})],
        [("", {"size": 16})],
        [("Fused CUDA kernel",
          {"size": 17, "bold": True, "color": NAVY, "font": SERIF})],
        [("·  everything on-GPU;  no global-memory round trips",
          {"size": 14})],
        [("·  early exit on partial convolution sums  (most compositions die cheap)",
          {"size": 14})],
        [("·  two-stage freeze kernel at deep levels — perturbation bound",
          {"size": 14})],
        [("   clears 90 %+ of frozen cases, avoiding warp divergence",
          {"size": 14})],
        [("·  10 × – 100 ×  speedup over the naive enumerator",
          {"size": 14, "bold": True, "color": INK})],
    ])

    s.shapes.add_picture(str(FIG / "fig_cascade_schematic.png"),
                         Inches(7.55), Inches(1.95), width=Inches(5.3))

    add_footer(s, idx, total, active_color=NAVY)


def slide_06_cascade_result(prs, idx, total):
    s = blank_slide(prs, rail_color=NAVY)
    slide_header(s, "TRACK I · RESULT",
                 "Cascade pushes the lower bound from 1.2802 to 1.4",
                 color=NAVY, eyebrow_color=NAVY)

    # Big editorial result block, left
    add_text(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.4),
             "CASCADE   ·   GPU", size=11, bold=True, color=NAVY, font=SANS)
    add_rule(s, Inches(0.6), Inches(2.3), Inches(1.8),
             color=NAVY, thickness=1.4)

    add_text(s, Inches(0.6), Inches(2.55), Inches(6.5), Inches(1.7),
             "C₁ₐ ≥ 1.4",
             size=96, bold=False, color=INK, font=SERIF, spacing=1.0)

    add_rich(s, Inches(0.6), Inches(4.55), Inches(6.0), Inches(1.8), [
        [("at  ", {"size": 18, "color": STONE}),
         ("(d, m) = (128, 20)",
          {"size": 18, "color": INK, "italic": True, "font": SERIF})],
        [("≈ 70 hours  on a 128-core machine",
          {"size": 16, "color": STONE, "italic": True})],
    ])

    add_rule(s, Inches(0.6), Inches(5.85), Inches(6.0),
             color=HAIRLINE, thickness=0.75)
    add_rich(s, Inches(0.6), Inches(5.95), Inches(6.0), Inches(1.2), [
        [("First significant lower-bound improvement since 2017.",
          {"size": 13, "color": INK})],
        [("Gap to upper bound  narrowed by ≈ 50 %  (0.22 → 0.10).",
          {"size": 13, "color": INK})],
    ])

    # Right column: trade-offs
    add_rule(s, Inches(7.2), Inches(1.95), Inches(0.4),
             color=NAVY, thickness=1.4)
    add_text(s, Inches(7.2), Inches(2.15), Inches(5.5), Inches(0.4),
             "TRADE-OFFS",
             size=11, bold=True, color=NAVY, font=SANS)

    add_rich(s, Inches(7.2), Inches(2.55), Inches(5.7), Inches(4.7), [
        [("What the cascade buys you",
          {"size": 16, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 10})],
        [("·  sharp in principle — exhaustive enumeration rules out every surviving point",
          {"size": 13})],
        [("·  scales to  ", {"size": 13}),
         ("d = 128", {"size": 13, "font": SERIF, "italic": True}),
         ("  via symmetry reduction + GPU pruning", {"size": 13})],
        [("", {"size": 16})],
        [("What it doesn't",
          {"size": 16, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 10})],
        [("·  exponential in  ", {"size": 13}),
         ("d", {"size": 13, "font": SERIF, "italic": True}),
         ("  in the worst case", {"size": 13})],
        [("·  correctness rides on many small inequalities — each needs a proof",
          {"size": 13})],
        [("·  motivated us to build a second, independent certifier …",
          {"size": 13, "italic": True, "color": STONE})],
    ])

    add_footer(s, idx, total, active_color=NAVY)


def slide_07_sdp_method(prs, idx, total):
    s = blank_slide(prs, rail_color=FOREST)
    slide_header(s, "TRACK II · LASSERRE SDP",
                 "Continuous → polynomial → semidefinite   (no enumeration)",
                 color=FOREST, eyebrow_color=FOREST)

    # Step 1
    add_rich(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.5), [
        [("1.  Reduce", {"size": 17, "bold": True, "color": FOREST,
                         "font": SERIF}),
         ("   the infinite-dimensional problem on  ",
          {"size": 16, "color": INK}),
         ("𝓕", {"size": 16, "italic": True, "font": SERIF}),
         ("  to a polynomial minimization on the simplex:",
          {"size": 16, "color": INK})],
    ])
    add_box(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(0.8),
            fill=HIGHLIGHT)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.6),
             "val(d)  =  minᵤ ∈ Δ_d    max_W    μᵀ M_W μ     ≤     C₁ₐ",
             size=22, bold=False, italic=True, color=INK, font=SERIF,
             align=PP_ALIGN.CENTER)

    # Step 2
    add_rich(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.5), [
        [("2.  Relax", {"size": 17, "bold": True, "color": FOREST,
                        "font": SERIF}),
         ("   via the Lasserre hierarchy with localizing matrices:",
          {"size": 16, "color": INK})],
    ])
    add_box(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(0.8),
            fill=HIGHLIGHT)
    add_text(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(0.6),
             "val⁽ᵏ⁾(d)    ≤    val⁽ᵏ⁺¹⁾(d)    ≤    val(d)",
             size=22, bold=False, italic=True, color=INK, font=SERIF,
             align=PP_ALIGN.CENTER)

    # Step 3
    add_rich(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(0.5), [
        [("3.  Certify", {"size": 17, "bold": True, "color": FOREST,
                          "font": SERIF}),
         ("   a dual-feasible point of the order-",
          {"size": 16, "color": INK}),
         ("k", {"size": 16, "italic": True, "font": SERIF}),
         ("  SDP whose objective exceeds the threshold.",
          {"size": 16, "color": INK})],
    ])

    # Scaling lever
    add_rule(s, Inches(0.6), Inches(6.05), Inches(12.1),
             color=HAIRLINE, thickness=0.75)
    add_rich(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(1.2), [
        [("Scaling lever",
          {"size": 14, "bold": True, "color": FOREST, "font": SERIF}),
         ("   —   correlative sparsity  (Waki–Kim–Kojima–Muramatsu 2006).",
          {"size": 14, "color": INK})],
        [("Chordal extension of the banded window graph yields  ",
          {"size": 13, "color": STONE}),
         ("O(d)", {"size": 13, "color": STONE, "font": SERIF,
                   "italic": True}),
         ("  overlapping cliques;  the dense  ", {"size": 13, "color": STONE}),
         ("binom(d+k, k)", {"size": 13, "color": STONE, "font": SERIF,
                            "italic": True}),
         ("  moment matrix is replaced by  ", {"size": 13, "color": STONE}),
         ("O(d)", {"size": 13, "color": STONE, "font": SERIF,
                   "italic": True}),
         ("  small blocks.", {"size": 13, "color": STONE})],
    ])

    add_footer(s, idx, total, active_color=FOREST)


def slide_08_sdp_result(prs, idx, total):
    s = blank_slide(prs, rail_color=FOREST)
    slide_header(s, "TRACK II · RESULT",
                 "A certified lower bound with full rigor",
                 color=FOREST, eyebrow_color=FOREST)

    s.shapes.add_picture(str(FIG / "fig_lasserre_ladder.png"),
                         Inches(0.4), Inches(1.85), width=Inches(7.5))

    # Right: big editorial result
    add_text(s, Inches(8.1), Inches(1.95), Inches(4.7), Inches(0.4),
             "CERTIFIED", size=11, bold=True, color=FOREST, font=SANS)
    add_rule(s, Inches(8.1), Inches(2.3), Inches(1.6),
             color=FOREST, thickness=1.4)
    add_text(s, Inches(8.1), Inches(2.5), Inches(5.0), Inches(1.6),
             "C₁ₐ ≥ 1.3",
             size=72, bold=False, color=INK, font=SERIF, spacing=1.0)
    add_rich(s, Inches(8.1), Inches(4.0), Inches(5.0), Inches(0.9), [
        [("at  ", {"size": 16, "color": STONE}),
         ("(d, k) = (16, 3)", {"size": 16, "color": INK, "italic": True,
                               "font": SERIF})],
        [("MOSEK  ·  < 10 minutes",
          {"size": 14, "color": STONE, "italic": True})],
    ])

    add_rule(s, Inches(8.1), Inches(5.1), Inches(4.7),
             color=HAIRLINE, thickness=0.75)
    add_rich(s, Inches(8.1), Inches(5.2), Inches(4.9), Inches(2.0), [
        [("What we record",
          {"size": 14, "bold": True, "color": FOREST, "font": SERIF})],
        [("", {"size": 6})],
        [("·  a single dual-feasible point  (y, Z, Λ)",
          {"size": 12, "color": INK})],
        [("·  post-processed via Jansson–Chaykin–Keil",
          {"size": 12, "color": INK})],
        [("   interval-arithmetic shift",
          {"size": 12, "color": INK})],
        [("·  auxiliary runs at  ", {"size": 12, "color": INK}),
         ("d ∈ {32, 64, 128}", {"size": 12, "color": INK, "italic": True,
                                 "font": SERIF})],
        [("   under correlative sparsity",
          {"size": 12, "color": INK})],
    ], spacing=1.3)

    add_footer(s, idx, total, active_color=FOREST)


def slide_09_matlab_setup(prs, idx, total):
    s = blank_slide(prs, rail_color=BRICK)
    slide_header(s, "TRACK III · AUDITING THE C&S MATLAB",
                 "The code the authors shared — what it does",
                 color=BRICK, eyebrow_color=BRICK)

    add_rich(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(2.6), [
        [("When we set out to reproduce the  ", {"size": 17, "color": INK}),
         ("C₁ₐ ≥ 1.2802", {"size": 17, "color": INK, "italic": True,
                            "font": SERIF}),
         ("  baseline we ran the MATLAB file", {"size": 17, "color": INK})],
        [("Cloninger and Steinerberger shared  (",
          {"size": 17, "color": INK}),
         ("original_baseline_matlab.m",
          {"size": 15, "color": INK, "font": MONO}),
         (").  The published C&S paper is ", {"size": 17, "color": INK}),
         ("correct", {"size": 17, "bold": True, "color": INK}),
         (" —",  {"size": 17, "color": INK})],
        [("it uses the fine height grid  ", {"size": 17, "color": INK}),
         ("B", {"size": 17, "italic": True, "font": SERIF}),
         ("n,m", {"size": 12, "italic": True, "font": SERIF}),
         ("  with  ", {"size": 17}),
         ("m = 50", {"size": 17, "italic": True, "font": SERIF}),
         (".", {"size": 17})],
        [("", {"size": 8})],
        [("The MATLAB file,  though,  parameterizes the search differently:",
          {"size": 17, "color": INK})],
    ])

    add_box(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.0),
            fill=HIGHLIGHT)
    add_rule(s, Inches(0.9), Inches(4.5), Inches(11.5),
             color=BRICK, thickness=1.3)
    add_rich(s, Inches(1.2), Inches(4.7), Inches(11.1), Inches(1.7), [
        [("In the paper",
          {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         (":    step-function  ", {"size": 15, "color": INK}),
         ("heights", {"size": 15, "italic": True, "font": SERIF}),
         ("   quantized in steps of   ", {"size": 15, "color": INK}),
         ("1/m = 0.02",
          {"size": 15, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 12})],
        [("In the code",
          {"size": 15, "bold": True, "color": BRICK, "font": SERIF}),
         (":     bin  ", {"size": 15, "color": INK}),
         ("masses", {"size": 15, "italic": True, "font": SERIF}),
         ("     quantized in steps of   ", {"size": 15, "color": INK}),
         ("0.02", {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         ("    →   heights quantized in steps of   ",
          {"size": 15, "color": INK}),
         ("2d × 0.02",
          {"size": 15, "bold": True, "color": BRICK, "font": SERIF})],
    ])

    add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
             "Different parameterization.  Different height step.  Same correction formula.",
             size=15, italic=True, color=STONE, align=PP_ALIGN.CENTER,
             font=SERIF)

    add_footer(s, idx, total, active_color=BRICK)


def slide_10_matlab_gap(prs, idx, total):
    s = blank_slide(prs, rail_color=BRICK)
    slide_header(s, "TRACK III · THE GAP",
                 "Lemma 3 needs the height step;  the code plugs in the mass step",
                 color=BRICK, eyebrow_color=BRICK)

    add_rich(s, Inches(0.6), Inches(1.8), Inches(6.4), Inches(5.2), [
        [("C&S Lemma 3  gives",
          {"size": 15, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 4})],
        [("(g ∗ g)(x)   ≤   (a ∗ a)(x)  +  2ε  +  ε²",
          {"size": 17, "color": INK, "font": SERIF, "italic": True})],
        [("where  ", {"size": 14}),
         ("ε", {"size": 14, "italic": True, "font": SERIF}),
         ("  is a bound on  ", {"size": 14}),
         ("‖g − a‖∞", {"size": 14, "font": SERIF, "italic": True}),
         ("   (a  ", {"size": 14}),
         ("height", {"size": 14, "italic": True, "font": SERIF}),
         ("   difference).", {"size": 14})],
        [("", {"size": 14})],
        [("In the MATLAB code",
          {"size": 15, "bold": True, "color": BRICK, "font": SERIF})],
        [("", {"size": 4})],
        [("ε = 0.02", {"size": 14, "italic": True, "font": SERIF}),
         ("  (mass step),  but heights  ", {"size": 14}),
         ("hᵢ = 2d · massᵢ",
          {"size": 14, "font": SERIF, "italic": True}),
         (",", {"size": 14})],
        [("so the true  ", {"size": 14}),
         ("‖g − a‖∞",
          {"size": 14, "font": SERIF, "italic": True}),
         ("  ≤  ", {"size": 14}),
         ("2d × 0.02", {"size": 14, "font": SERIF, "italic": True}),
         (",   not  ", {"size": 14}),
         ("0.02",
          {"size": 14, "font": SERIF, "italic": True}),
         (".", {"size": 14})],
        [("", {"size": 16})],
        [("The code's pruning budget is too small by a factor of  ",
          {"size": 15, "color": INK})],
        [("≈ 2d.",
          {"size": 20, "bold": True, "color": BRICK, "font": SERIF})],
    ])

    s.shapes.add_picture(str(FIG / "fig_correction_gap.png"),
                         Inches(7.1), Inches(1.9), width=Inches(5.9))

    add_footer(s, idx, total, active_color=BRICK)


def slide_11_smoking_gun(prs, idx, total):
    s = blank_slide(prs, rail_color=BRICK)
    slide_header(s, "TRACK III · THE CONSEQUENCE",
                 'The code would "prove away" a known construction',
                 color=BRICK, eyebrow_color=BRICK)

    add_rich(s, Inches(0.6), Inches(1.8), Inches(6.3), Inches(5.2), [
        [("The check",
          {"size": 16, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 4})],
        [("Matolcsi–Vinuesa  (2010)  exhibit a continuous  ",
          {"size": 14}),
         ("f*", {"size": 14, "italic": True, "font": SERIF}),
         ("  with", {"size": 14})],
        [("‖f* ∗ f*‖∞  ≤  1.5029",
          {"size": 15, "color": INK, "italic": True, "font": SERIF}),
         (".", {"size": 14})],
        [("So  ", {"size": 14}),
         ("C₁ₐ ≤ 1.5029",
          {"size": 14, "italic": True, "font": SERIF}),
         (".", {"size": 14, "italic": True})],
        [("", {"size": 14})],
        [("Point the MATLAB cascade at the threshold  ",
          {"size": 14}),
         ("1.51", {"size": 14, "font": SERIF, "italic": True}),
         (":", {"size": 14})],
        [("·  interior-window budget stays at  ≈ 0.04  across levels",
          {"size": 13})],
        [("·  the step-function approximation of  ",
          {"size": 13}),
         ("f*", {"size": 13, "italic": True, "font": SERIF}),
         ("  has self-convolution peaks", {"size": 13})],
        [("   that exceed 1.55 on narrow windows at moderate  ",
          {"size": 13}),
         ("d", {"size": 13, "italic": True, "font": SERIF})],
        [("·  the code therefore prunes every descendant of  ",
          {"size": 13}),
         ("f*", {"size": 13, "italic": True, "font": SERIF})],
        [("·  exhausted survivors  →  claim:  ",
          {"size": 13}),
         ("C₁ₐ ≥ 1.51", {"size": 13, "font": SERIF, "italic": True,
                          "bold": True})],
        [("", {"size": 14})],
        [("But  ", {"size": 15, "bold": True, "color": BRICK,
                     "font": SERIF}),
         ("1.51 > 1.5029 ≥ C₁ₐ",
          {"size": 15, "bold": True, "color": BRICK, "font": SERIF,
           "italic": True}),
         (".    Contradiction.",
          {"size": 15, "bold": True, "color": BRICK, "font": SERIF})],
    ])

    s.shapes.add_picture(str(FIG / "fig_effective_m.png"),
                         Inches(7.3), Inches(1.9), width=Inches(5.7))
    add_rule(s, Inches(7.3), Inches(5.8), Inches(5.7),
             color=HAIRLINE, thickness=0.75)
    add_text(s, Inches(7.3), Inches(5.95), Inches(5.7), Inches(1.1),
             "The effective fine-grid  m  in C&S units is  1 / (2dε)  — "
             "it drops below 1 by  d ≈ 25, at which point Lemma 3 says nothing at all.",
             size=11, color=STONE, italic=True, font=SANS)

    add_footer(s, idx, total, active_color=BRICK)


def slide_12_what_this_means(prs, idx, total):
    s = blank_slide(prs, rail_color=BRICK)
    slide_header(s, "TRACK III · WHAT WE ARE SAYING",
                 "The paper stands;  the shared MATLAB file does not prove it",
                 color=BRICK, eyebrow_color=BRICK)

    # Left: the paper is fine
    add_text(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.4),
             "THE PAPER IS FINE",
             size=11, bold=True, color=INK, font=SANS)
    add_rule(s, Inches(0.6), Inches(2.2), Inches(1.6),
             color=INK, thickness=1.3)
    add_rich(s, Inches(0.6), Inches(2.4), Inches(5.9), Inches(4.4), [
        [("C&S 2017 work on the fine height grid  ",
          {"size": 14, "color": INK}),
         ("B", {"size": 14, "italic": True, "font": SERIF}),
         ("n,m", {"size": 10, "italic": True, "font": SERIF}),
         (",", {"size": 14, "color": INK})],
        [("where Lemma 3's  ", {"size": 14, "color": INK}),
         ("2/m + 1/m²", {"size": 14, "italic": True, "font": SERIF}),
         ("   correction is legitimate.", {"size": 14, "color": INK})],
        [("", {"size": 10})],
        [("The published bound  ", {"size": 14, "color": INK}),
         ("1.2802", {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         ("  is not in doubt.", {"size": 14, "color": INK})],
        [("", {"size": 14})],
        [("The MATLAB artifact is a different animal:",
          {"size": 14, "color": INK})],
        [("it enumerates a coarser mass grid, but ports the paper's",
          {"size": 14, "color": INK})],
        [("correction formula verbatim.  That combination",
          {"size": 14, "color": INK})],
        [("cannot certify the paper's bound, and as shown would",
          {"size": 14, "color": INK})],
        [("certify bounds that contradict the MV upper bound.",
          {"size": 14, "color": INK})],
    ])

    # Right: two hypotheses
    add_text(s, Inches(6.95), Inches(1.85), Inches(6.0), Inches(0.4),
             "TWO LIVE HYPOTHESES",
             size=11, bold=True, color=BRICK, font=SANS)
    add_rule(s, Inches(6.95), Inches(2.2), Inches(1.6),
             color=BRICK, thickness=1.3)
    add_rich(s, Inches(6.95), Inches(2.4), Inches(5.9), Inches(4.4), [
        [("(a)  They sent us an early prototype.",
          {"size": 14, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 4})],
        [("The published bound was run on a separate fine-grid",
          {"size": 13, "color": INK})],
        [("enumerator;  the coarse-mass file is a speed-oriented",
          {"size": 13, "color": INK})],
        [("variant never meant to be the proof vehicle.",
          {"size": 13, "color": INK})],
        [("", {"size": 14})],
        [("(b)  A derivation we're missing.",
          {"size": 14, "bold": True, "color": INK, "font": SERIF})],
        [("", {"size": 4})],
        [("A mass-space bound that avoids the  ",
          {"size": 13, "color": INK}),
         ("2d", {"size": 13, "italic": True, "font": SERIF}),
         ("  factor — ", {"size": 13, "color": INK})],
        [("unlikely, but worth asking about before we conclude.",
          {"size": 13, "color": INK})],
        [("", {"size": 14})],
        [("Either way,  this is a narrow claim about one file.",
          {"size": 13, "italic": True, "color": STONE, "font": SERIF})],
    ])

    add_rule(s, Inches(0.6), Inches(6.7), Inches(12.1),
             color=HAIRLINE, thickness=0.75)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.45),
             "We fixed the analogous bug in our own Python port on 2026-04-07 — "
             "switched to the fine grid, which is what makes the cascade's 1.4 sound.",
             size=12, italic=True, color=STONE,
             align=PP_ALIGN.CENTER, font=SANS)

    add_footer(s, idx, total, active_color=BRICK)


def slide_13_summary(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "§3 · Takeaways",
                 "Where C₁ₐ stands after this project")

    s.shapes.add_picture(str(FIG / "fig_bounds_number_line.png"),
                         Inches(1.9), Inches(1.7), width=Inches(9.5))

    add_rule(s, Inches(0.6), Inches(4.4), Inches(12.1),
             color=HAIRLINE, thickness=0.75)

    add_text(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(0.4),
             "THREE DELIVERABLES", size=11, bold=True, color=STONE,
             font=SANS)

    add_rich(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(2.3), [
        [("I.   ", {"size": 15, "bold": True, "color": NAVY,
                    "font": SERIF}),
         ("Cascade on GPU",
          {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         ("   —   ", {"size": 14, "color": STONE}),
         ("C₁ₐ ≥ 1.4",
          {"size": 15, "bold": True, "color": NAVY, "font": SERIF,
           "italic": True}),
         ("   via novel pruning bounds + fused CUDA kernels.",
          {"size": 14, "color": INK})],
        [("II.  ", {"size": 15, "bold": True, "color": FOREST,
                    "font": SERIF}),
         ("Lasserre SDP certificate",
          {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         ("   —   ", {"size": 14, "color": STONE}),
         ("C₁ₐ ≥ 1.3  (certified)",
          {"size": 15, "bold": True, "color": FOREST, "font": SERIF,
           "italic": True}),
         ("   at  ", {"size": 14, "color": INK}),
         ("(d, k) = (16, 3)", {"size": 14, "italic": True, "font": SERIF}),
         ("   in under 10 minutes.", {"size": 14, "color": INK})],
        [("III. ", {"size": 15, "bold": True, "color": BRICK,
                    "font": SERIF}),
         ("MATLAB audit",
          {"size": 15, "bold": True, "color": INK, "font": SERIF}),
         ("   —   the shared C&S artifact confuses the mass step for the height step; ",
          {"size": 14, "color": INK}),
         ("the published paper is unaffected.",
          {"size": 14, "color": INK, "italic": True})],
        [("", {"size": 12})],
        [("Next:   formal write-ups for both proofs in  ",
          {"size": 13, "color": STONE}),
         ("proof/", {"size": 13, "color": STONE, "font": MONO}),
         (",   push cascade to  d ≥ 256,   close the gap toward 1.5029.",
          {"size": 13, "color": STONE, "italic": True})],
    ], spacing=1.45)

    add_footer(s, idx, total, active_color=INK)


# ============================================================
# BUILD
# ============================================================

def build():
    prs = new_prs()
    slide_01_title(prs)
    TOTAL = 12

    slide_02_problem(prs, 2, TOTAL)
    slide_03_baseline_and_tracks(prs, 3, TOTAL)
    slide_05_cascade_method(prs, 4, TOTAL)
    slide_06_cascade_result(prs, 5, TOTAL)
    slide_07_sdp_method(prs, 6, TOTAL)
    slide_08_sdp_result(prs, 7, TOTAL)
    slide_09_matlab_setup(prs, 8, TOTAL)
    slide_10_matlab_gap(prs, 9, TOTAL)
    slide_11_smoking_gun(prs, 10, TOTAL)
    slide_12_what_this_means(prs, 11, TOTAL)
    slide_13_summary(prs, 12, TOTAL)

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    build()
